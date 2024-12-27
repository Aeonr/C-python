from pptx.util import Inches
from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 添加表格到幻灯片
table = slide.shapes.add_table(rows=4, cols=3, left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(5))
table.table.style = 'Table Grid'
table.table.rows[0].height = Inches(4)
table.table.rows[1].height = Inches(0.5)
table.table.rows[2].height = Inches(0.5)
table.table.rows[3].height = Inches(0.5)

# 在幻灯片中添加图像
img_path = r"C:\Users\CYH10\Desktop\微信图片_20241115170539.png"
slide.shapes.add_picture(img_path, Inches(0.4), Inches(1), Inches(4), Inches(4))

prs.save('python_course.pptx')

